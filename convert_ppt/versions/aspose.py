import os
import re
from pptx import Presentation
from docx import Document
from docx.shared import Inches
import aspose.slides as slides
import aspose.pydrawing as drawing

# text clean-up
def sanitize_text(text):
    return re.sub(r'[^\x20-\x7E]', '', text)

# func to create img from slide
def generate_thumbnails(presentation_path, output_folder):
    pres = slides.Presentation(presentation_path)
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    scaleX, scaleY = 2.0, 2.0 #give size of the img
    for index in range(pres.slides.length):
        slide = pres.slides[index]
        slide.get_thumbnail(scaleX, scaleY).save(os.path.join(output_folder, f"slide_{index}.jpg"), drawing.imaging.ImageFormat.jpeg)

# slide processing
def process_slide(slide, doc, slide_index):
    title = sanitize_text(slide.shapes.title.text if slide.shapes.title else "")
    body = sanitize_text("\n".join([shape.text for shape in slide.shapes if shape.has_text_frame]))

    doc.add_heading(title, level=1)

    image_path = f"result/slide_{slide_index}.jpg"
    if os.path.exists(image_path):
        doc.add_picture(image_path, width=Inches(6)) #give size to display in word doc

    if body:
        doc.add_paragraph(body)

    return doc
# to handle presentaion process
def process_presentation(input_pptx, output_docx):
    try:
        generate_thumbnails(input_pptx, "result")
        
        doc = Document()
        presentation = Presentation(input_pptx)

        for index, slide in enumerate(presentation.slides):
            process_slide(slide, doc, index)
            if index < len(presentation.slides) - 1:
                doc.add_page_break()

        doc.save(output_docx)   # output document
        print("Document successfully created:", output_docx)

    except Exception as e:
        print("An error occurred:", e)


process_presentation("Unit1.pptx", "output.docx")
